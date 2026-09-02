from __future__ import annotations

import base64, io, json, os, re, shutil, subprocess, tempfile, wave, zipfile
from pathlib import Path

import requests
from docx import Document
from fastapi import APIRouter, File, Form, HTTPException, UploadFile
from fastapi.responses import FileResponse
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, PatternFill
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Image as PdfImage, Paragraph, SimpleDocTemplate, Spacer

router = APIRouter(prefix="/intelligence", tags=["video intelligence"])
JOBS = Path(__file__).resolve().parent / "outputs" / "intelligence"
JOBS.mkdir(parents=True, exist_ok=True)


def ffmpeg():
    root = Path(__file__).resolve().parent.parent
    bundled = root / "node_modules/ffmpeg-static" / ("ffmpeg.exe" if os.name == "nt" else "ffmpeg")
    return shutil.which("ffmpeg") or (str(bundled) if bundled.exists() else None)


def stamp(seconds):
    ms = int(round(seconds * 1000)); h, ms = divmod(ms, 3600000); m, ms = divmod(ms, 60000); s, ms = divmod(ms, 1000)
    return f"{h:02}:{m:02}:{s:02},{ms:03}"


def transcribe(path: Path, language=None, task="transcribe"):
    try:
        from faster_whisper import WhisperModel
    except ImportError:
        raise HTTPException(503, "faster-whisper is not installed")
    model = WhisperModel("small", device="cpu", compute_type="int8")
    segments, info = model.transcribe(str(path), vad_filter=True, language=language or None, word_timestamps=True, task=task)
    rows = [{"start": float(x.start), "end": float(x.end), "text": x.text.strip(),"speaker":"Speaker 1","confidence":round(max(0,min(1,1-float(getattr(x,"avg_logprob",-1))/-5)),2),"words":[{"start":float(w.start),"end":float(w.end),"word":w.word,"confidence":round(float(getattr(w,"probability",0)),2)} for w in (x.words or [])]} for x in segments]
    return rows, getattr(info, "language", language or "unknown")


def extract_frames(video: Path, folder: Path, interval: int, scene_aware=True):
    exe = ffmpeg()
    if not exe: return []
    folder.mkdir(exist_ok=True)
    pattern = folder / "frame-%04d.jpg"
    selector="select='gt(scene,0.28)',scale=1280:-2" if scene_aware else f"fps=1/{max(3, interval)},scale=1280:-2"
    subprocess.run([exe, "-y", "-i", str(video), "-vf", selector, "-vsync", "vfr", "-q:v", "3", str(pattern)], capture_output=True, check=True)
    return sorted(folder.glob("frame-*.jpg"))[:80]


def ocr_image(path: Path):
    try:
        import pytesseract
        return pytesseract.image_to_string(Image.open(path)).strip()
    except Exception:
        return ""


def vision_note(path: Path):
    encoded = base64.b64encode(path.read_bytes()).decode()
    prompt = "Describe this video frame for study notes. Extract visible text, charts, tables, diagrams, equations and formulas faithfully. Return concise plain text. Do not invent unreadable content."
    for model in ("qwen2.5vl:3b", "llava:7b"):
        try:
            response = requests.post("http://127.0.0.1:11434/api/generate", json={"model": model, "prompt": prompt, "images": [encoded], "stream": False}, timeout=120)
            if response.ok: return response.json().get("response", "").strip()
        except Exception: pass
    return ""


def summarize(transcript: str, visual_notes: list[dict], mode="standard", purpose="lecture"):
    prompt = f"Create {mode} accurate {purpose} notes. Return JSON with title, summary, key_points, chapters (each with time and title), terms, formulas (each with latex, mathml, source_time, confidence), tables (headers, rows, source_time, confidence), flashcards, quiz, mind_map, action_items, decisions, owners, deadlines. Every claim needs source_time. Preserve uncertainty; never invent missing detail.\n" + json.dumps({"transcript": transcript[:50000], "visuals": visual_notes[:40]})
    try:
        r = requests.post("http://127.0.0.1:11434/api/generate", json={"model":"qwen2.5:3b","prompt":prompt,"stream":False,"format":"json"}, timeout=180)
        r.raise_for_status(); return json.loads(r.json()["response"])
    except Exception:
        sentences = [x.strip() for x in re.split(r"(?<=[.!?])\s+", transcript) if x.strip()]
        return {"title":"Video notes","summary":" ".join(sentences[:5]),"key_points":sentences[:12],"chapters":[],"terms":[],"formulas":[],"tables":[],"action_items":[]}


@router.post("/analyze")
async def analyze(file: UploadFile = File(...), language: str = Form(""), frame_interval: int = Form(20), note_mode:str=Form("standard"), purpose:str=Form("lecture"), target_language:str=Form(""), scene_aware:bool=Form(True), translate_to_english:bool=Form(False)):
    work = Path(tempfile.mkdtemp(prefix="video-notes-")); source = work / (file.filename or "video.mp4")
    source.write_bytes(await file.read())
    try:
        segments, detected = transcribe(source, language)
        frames = extract_frames(source, work / "frames", frame_interval, scene_aware)
        visuals=[]
        for index, frame in enumerate(frames):
            ocr=ocr_image(frame); description=vision_note(frame)
            if ocr or description:
                visuals.append({"time":index*max(3,frame_interval),"ocr":ocr,"description":description,"confidence":0.9 if ocr else 0.55,"image":"data:image/jpeg;base64,"+base64.b64encode(frame.read_bytes()).decode()})
        transcript=" ".join(x["text"] for x in segments)
        english_segments=[]
        if translate_to_english:
            english_segments=segments if detected.lower().startswith("en") else transcribe(source,language or None,task="translate")[0]
        notes=summarize(transcript,visuals,note_mode,purpose)
        translation=""
        if target_language:
            try:
                p=f"Translate this transcript into {target_language}, preserving timestamps and speaker labels. Return plain text only.\n"+"\n".join(f"[{stamp(x['start'])[:-4]}] {x['speaker']}: {x['text']}" for x in segments)
                r=requests.post("http://127.0.0.1:11434/api/generate",json={"model":"qwen2.5:3b","prompt":p,"stream":False},timeout=180);translation=r.json().get("response","")
            except Exception: pass
        srt="\n\n".join(f"{i+1}\n{stamp(x['start'])} --> {stamp(x['end'])}\n{x['text']}" for i,x in enumerate(segments))+"\n"
        english_srt="\n\n".join(f"{i+1}\n{stamp(x['start'])} --> {stamp(x['end'])}\n{x['text']}" for i,x in enumerate(english_segments))+("\n" if english_segments else "")
        return {"filename":file.filename,"language":detected,"translation":translation,"transcript":transcript,"segments":segments,"srt":srt,"english_segments":english_segments,"english_srt":english_srt,"notes":notes,"visuals":visuals,"settings":{"mode":note_mode,"purpose":purpose,"scene_aware":scene_aware,"translated_to_english":translate_to_english}}
    finally: shutil.rmtree(work,ignore_errors=True)


def add_doc_content(doc, data, images=True):
    notes=data.get("notes",{}); doc.add_heading(notes.get("title") or "Video notes",0); doc.add_paragraph(notes.get("summary", ""))
    for key,label in (("key_points","Key points"),("chapters","Chapters"),("terms","Terms"),("formulas","Math and formulas"),("action_items","Action items")):
        values=notes.get(key) or []
        if values:
            doc.add_heading(label,1)
            for value in values: doc.add_paragraph(str(value),style="List Bullet")
    if images:
        for item in data.get("visuals",[])[:20]:
            try:
                raw=base64.b64decode(item["image"].split(",",1)[1]); doc.add_picture(io.BytesIO(raw),width=Inches(5.8)); doc.add_paragraph(f"Frame at {item['time']}s - {item.get('ocr') or item.get('description','')}")
            except Exception: pass
    doc.add_heading("Transcript",1)
    for s in data.get("segments",[]): doc.add_paragraph(f"[{stamp(s['start'])[:-4]}] {s['text']}")


@router.post("/export")
async def export_notes(payload: str = Form(...), format: str = Form("docx")):
    data=json.loads(payload); fmt=format.lower(); out=JOBS/f"video-notes.{fmt}"
    if fmt=="srt": out.write_text(data.get("srt",""),encoding="utf-8")
    elif fmt=="vtt": out.write_text("WEBVTT\n\n"+data.get("srt","").replace(",","."),encoding="utf-8")
    elif fmt=="english-srt":
        out=JOBS/"video-notes.english.srt";out.write_text(data.get("english_srt",""),encoding="utf-8")
    elif fmt=="english-vtt":
        out=JOBS/"video-notes.english.vtt";out.write_text("WEBVTT\n\n"+data.get("english_srt","").replace(",","."),encoding="utf-8")
    elif fmt=="docx":
        doc=Document(); add_doc_content(doc,data); doc.save(out)
    elif fmt=="pptx":
        prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
        notes=data.get("notes",{}); sections=[("Summary",notes.get("summary", "")),("Key points","\n".join(map(str,notes.get("key_points",[])))),("Math and formulas","\n".join(map(str,notes.get("formulas",[])))),("Transcript",data.get("transcript","")[:3500])]
        for title,body in sections:
            slide=prs.slides.add_slide(prs.slide_layouts[5]); slide.shapes.title.text=title
            box=slide.shapes.add_textbox(Inches(.8),Inches(1.5),Inches(11.7),Inches(5.2)); box.text_frame.text=body
            for p in box.text_frame.paragraphs:p.font.size=Pt(20)
        prs.save(out)
    elif fmt=="xlsx":
        wb=Workbook(); ws=wb.active; ws.title="Transcript"; ws.append(["Start","End","Text"])
        for s in data.get("segments",[]):ws.append([s["start"],s["end"],s["text"]])
        for cell in ws[1]:cell.font=Font(bold=True,color="FFFFFF");cell.fill=PatternFill("solid",fgColor="6657E9")
        ws.column_dimensions["A"].width=12;ws.column_dimensions["B"].width=12;ws.column_dimensions["C"].width=90
        visual=wb.create_sheet("Visual extraction");visual.append(["Time","OCR text","Description"])
        for x in data.get("visuals",[]):visual.append([x["time"],x.get("ocr",""),x.get("description","")])
        for index, table in enumerate(data.get("notes",{}).get("tables",[]),1):
            sheet=wb.create_sheet(f"Table {index}")
            if isinstance(table,dict):
                sheet.append(list(map(str,table.get("headers",[]))))
                for row in table.get("rows",[]):sheet.append(list(row) if isinstance(row,list) else [str(row)])
                sheet["A1"].comment=None
        wb.save(out)
    elif fmt=="pdf":
        styles=getSampleStyleSheet(); story=[Paragraph(data.get("notes",{}).get("title","Video notes"),styles["Title"]),Spacer(1,12),Paragraph(data.get("notes",{}).get("summary",""),styles["BodyText"])]
        for point in data.get("notes",{}).get("key_points",[]):story.extend([Spacer(1,7),Paragraph("- "+str(point),styles["BodyText"])])
        story.extend([Spacer(1,18),Paragraph("Transcript",styles["Heading1"]),Paragraph(data.get("transcript","").replace("\n","<br/>"),styles["BodyText"])])
        SimpleDocTemplate(str(out),pagesize=letter,rightMargin=54,leftMargin=54).build(story)
    else: raise HTTPException(400,"Supported formats: srt, vtt, docx, pptx, pdf, xlsx")
    return FileResponse(out,filename=out.name)


@router.post("/assets")
async def export_assets(payload:str=Form(...)):
    data=json.loads(payload);out=JOBS/"video-extracted-assets.zip"
    with zipfile.ZipFile(out,"w",zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("transcript.txt",data.get("transcript",""));archive.writestr("subtitles.srt",data.get("srt",""));archive.writestr("notes.json",json.dumps(data.get("notes",{}),indent=2))
        for i,item in enumerate(data.get("visuals",[]),1):
            try:archive.writestr(f"frames/frame-{i:04}.jpg",base64.b64decode(item["image"].split(",",1)[1]))
            except Exception:pass
    return FileResponse(out,filename=out.name)


@router.post("/narration")
async def narration(text:str=Form(...),format:str=Form("mp3"),rate:float=Form(1.0)):
    if not text.strip():raise HTTPException(400,"Add text to narrate")
    work=Path(tempfile.mkdtemp(prefix="narration-"));wav=work/"voice.wav";out=JOBS/f"narration.{format}"
    try:
        if os.name=="nt":
            safe=text[:100000].replace("'","''");target=str(wav).replace("'","''");speed=max(-10,min(10,round((rate-1)*10)))
            script=f"Add-Type -AssemblyName System.Speech;$s=New-Object System.Speech.Synthesis.SpeechSynthesizer;$s.Rate={speed};$s.SetOutputToWaveFile('{target}');$s.Speak('{safe}');$s.Dispose()"
            subprocess.run(["powershell","-NoProfile","-Command",script],check=True,capture_output=True)
        elif shutil.which("piper"):
            subprocess.run(["piper","--output_file",str(wav)],input=text.encode(),check=True)
        else:raise HTTPException(503,"Install a system speech voice or Piper")
        if format=="wav":shutil.copy2(wav,out)
        elif format=="mp3":
            if not ffmpeg():raise HTTPException(503,"FFmpeg is required")
            subprocess.run([ffmpeg(),"-y","-i",str(wav),"-codec:a","libmp3lame","-q:a","2",str(out)],check=True,capture_output=True)
        else:raise HTTPException(400,"Use mp3 or wav")
        return FileResponse(out,filename=out.name)
    finally:shutil.rmtree(work,ignore_errors=True)


def read_text(path: Path):
    suffix=path.suffix.lower()
    if suffix in (".txt",".md",".csv",".json",".srt",".vtt",".html"): return path.read_text(encoding="utf-8",errors="ignore")
    if suffix==".pdf":
        from pypdf import PdfReader
        return "\n".join(page.extract_text() or "" for page in PdfReader(path).pages)
    if suffix==".docx": return "\n".join(p.text for p in Document(path).paragraphs)
    if suffix==".pptx": return "\n".join(shape.text for slide in Presentation(path).slides for shape in slide.shapes if hasattr(shape,"text"))
    if suffix in (".xlsx",".xlsm"):
        wb=load_workbook(path,read_only=True,data_only=True); return "\n".join(" | ".join(str(v or "") for v in row) for ws in wb for row in ws.iter_rows(values_only=True))
    if suffix in (".png",".jpg",".jpeg",".webp",".tif",".tiff"): return ocr_image(path)
    raise HTTPException(415,"Unsupported file. Use PDF, DOCX, PPTX, XLSX, text, subtitle, HTML, JSON, CSV, or an image.")


@router.post("/read-file")
async def read_file(file: UploadFile=File(...), rate: float=Form(1.0)):
    work=Path(tempfile.mkdtemp(prefix="read-file-")); path=work/(file.filename or "document.txt"); path.write_bytes(await file.read())
    try:
        text=read_text(path).strip()
        if not text: raise HTTPException(422,"No readable text was found")
        return {"filename":file.filename,"text":text[:200000],"characters":len(text)}
    finally: shutil.rmtree(work,ignore_errors=True)
